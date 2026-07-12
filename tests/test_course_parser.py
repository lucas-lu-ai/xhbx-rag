import json
from pathlib import Path

import pytest

from xhbx_rag.chunk_io import load_chunks_jsonl
from xhbx_rag.course_enrichment import CourseEnrichment
from xhbx_rag.course_parser import CourseFileParseError, parse_course_dir


def _write_pptx(path: Path, slide_texts: list[str]) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    for text in slide_texts:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        textbox.text_frame.text = text
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(path))


class _FakeEnrichmentAgent:
    def __init__(self) -> None:
        self.calls: list[str] = []

    def enrich(self, course_name: str, course_series: str, sample_text: str) -> CourseEnrichment:
        self.calls.append(course_name)
        return CourseEnrichment(summary=f"{course_name}的课程摘要", audience="新人")


class _FailingEnrichmentAgent:
    def enrich(self, course_name: str, course_series: str, sample_text: str) -> CourseEnrichment:
        raise RuntimeError("模型不可用")


def test_parse_course_dir_scans_recursively_and_writes_outputs(tmp_path) -> None:
    course_dir = tmp_path / "培训数据"
    _write_pptx(
        course_dir / "新人培训" / "促成课.pptx",
        ["促成及异议处理培训课程，讲解促成动作要点与异议处理方法，" * 5],
    )
    (course_dir / "新人培训" / "~$促成课.pptx").write_bytes(b"junk")
    (course_dir / "新人培训" / "视频.mp4").write_bytes(b"junk")
    (course_dir / "新人培训" / "老教材.doc").write_bytes(b"junk")
    out_dir = tmp_path / "out"

    report = parse_course_dir(course_dir, out_dir, enrichment_agent=None)

    chunks = load_chunks_jsonl(out_dir / "chunks.jsonl")
    assert chunks
    assert all(chunk.chunk_type == "training_course" for chunk in chunks)
    assert chunks[0].metadata["course_series"] == "新人培训"

    report_data = json.loads((out_dir / "parse_report.json").read_text(encoding="utf-8"))
    assert report_data["counts"]["files_parsed"] == 1
    assert report_data["counts"]["chunks"] == len(chunks)
    skipped = {Path(item).name for item in report_data["skipped_files"]}
    assert skipped == {"~$促成课.pptx", "视频.mp4", "老教材.doc"}
    assert report.counts["files_parsed"] == 1


def test_parse_course_dir_isolates_single_file_failure(tmp_path) -> None:
    course_dir = tmp_path / "培训数据"
    _write_pptx(
        course_dir / "好课件.pptx",
        ["这是能正常解析的课件内容，讲解销售流程与客户经营方法，" * 5],
    )
    (course_dir / "坏课件.pptx").write_bytes(b"not a real pptx")
    out_dir = tmp_path / "out"

    report = parse_course_dir(course_dir, out_dir, enrichment_agent=None)

    assert report.counts["files_parsed"] == 1
    assert report.counts["files_failed"] == 1
    assert "坏课件.pptx" in report.failed_files[0]["path"]
    chunks = load_chunks_jsonl(out_dir / "chunks.jsonl")
    assert any("好课件" in chunk.metadata["course_name"] for chunk in chunks)


def test_parse_course_dir_fail_fast_raises_on_supported_file_failure(tmp_path) -> None:
    course_dir = tmp_path / "培训数据"
    _write_pptx(course_dir / "好课件.pptx", ["销售流程" * 50])
    (course_dir / "坏课件.pptx").write_bytes(b"broken")
    out_dir = tmp_path / "out"
    out_dir.mkdir()
    (out_dir / "chunks.jsonl").write_text("过期产物\n", encoding="utf-8")
    events: list[tuple[str, str, int]] = []

    with pytest.raises(CourseFileParseError, match="坏课件.pptx"):
        parse_course_dir(
            course_dir,
            out_dir,
            fail_fast=True,
            on_file=lambda path, status, count: events.append((path, status, count)),
        )

    assert events == [("坏课件.pptx", "failed", 0)]
    assert not (out_dir / "chunks.jsonl").exists()


def test_parse_course_dir_default_still_isolates_file_failure(tmp_path) -> None:
    course_dir = tmp_path / "培训数据"
    _write_pptx(course_dir / "好课件.pptx", ["销售流程" * 50])
    (course_dir / "坏课件.pptx").write_bytes(b"broken")

    report = parse_course_dir(course_dir, tmp_path / "out")

    assert report.counts["files_parsed"] == 1
    assert report.counts["files_failed"] == 1


def test_parse_course_dir_fail_fast_keeps_enrichment_failure_as_warning(tmp_path) -> None:
    course_dir = tmp_path / "培训数据"
    _write_pptx(course_dir / "促成课.pptx", ["促成动作要点" * 50])
    events: list[tuple[str, str, int]] = []

    report = parse_course_dir(
        course_dir,
        tmp_path / "out",
        enrichment_agent=_FailingEnrichmentAgent(),
        fail_fast=True,
        on_file=lambda path, status, count: events.append((path, status, count)),
    )

    assert report.enrich_failures and "促成课" in report.enrich_failures[0]
    assert events and events[0][0:2] == ("促成课.pptx", "parsed")
    assert events[0][2] > 0


def test_parse_course_dir_records_duplicate_text_hash(tmp_path) -> None:
    course_dir = tmp_path / "培训数据"
    body = "同一份课件内容被复制了两份，讲解客户约访的注意事项，" * 5
    (course_dir / "a").mkdir(parents=True)
    (course_dir / "a" / "课件.txt").write_text(body, encoding="utf-8")
    (course_dir / "a" / "课件(1).txt").write_text(body, encoding="utf-8")
    out_dir = tmp_path / "out"

    report = parse_course_dir(course_dir, out_dir, enrichment_agent=None)

    assert report.counts["duplicate_text_hashes"] >= 1


def test_parse_course_dir_applies_enrichment(tmp_path) -> None:
    course_dir = tmp_path / "培训数据"
    _write_pptx(
        course_dir / "促成课.pptx",
        ["促成动作要点讲解，帮助新人掌握促成与异议处理，" * 5],
    )
    out_dir = tmp_path / "out"
    agent = _FakeEnrichmentAgent()

    parse_course_dir(course_dir, out_dir, enrichment_agent=agent)

    assert agent.calls == ["促成课"]
    chunks = load_chunks_jsonl(out_dir / "chunks.jsonl")
    overview = next(chunk for chunk in chunks if chunk.chunk_id.endswith("__overview"))
    assert overview.metadata["summary"] == "促成课的课程摘要"


def test_parse_course_dir_degrades_when_enrichment_fails(tmp_path) -> None:
    course_dir = tmp_path / "培训数据"
    _write_pptx(
        course_dir / "促成课.pptx",
        ["促成动作要点讲解，帮助新人掌握促成与异议处理，" * 5],
    )
    out_dir = tmp_path / "out"

    report = parse_course_dir(course_dir, out_dir, enrichment_agent=_FailingEnrichmentAgent())

    assert report.enrich_failures == ["促成课: 课程增值服务不可用"]
    assert "模型不可用" not in report.enrich_failures[0]
    chunks = load_chunks_jsonl(out_dir / "chunks.jsonl")
    overview = next(chunk for chunk in chunks if chunk.chunk_id.endswith("__overview"))
    assert overview.metadata["summary"] == ""
