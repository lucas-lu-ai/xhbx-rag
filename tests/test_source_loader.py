from io import BytesIO

from PIL import Image

from xhbx_rag.models import EvidenceRef
from xhbx_rag.source_loader import (
    enrich_evidence_ref_location,
    load_case_sections,
    parse_source_file,
)


def _png_bytes() -> bytes:
    buffer = BytesIO()
    Image.new("RGB", (8, 8), "red").save(buffer, format="PNG")
    return buffer.getvalue()


def test_parse_txt_source_keeps_line_locator(tmp_path) -> None:
    source = tmp_path / "第1节.track-0.txt"
    source.write_text(
        "开场白\n客户说每年不能超过80万\n销售回应可以做预算释放\n",
        encoding="utf-8",
    )

    parsed = parse_source_file(source)

    assert parsed.source_type == "txt"
    assert parsed.filename == "第1节.track-0.txt"
    assert "客户说每年不能超过80万" in parsed.text

    ref = EvidenceRef(
        filename="第1节.track-0.txt",
        quote="客户说每年不能超过80万",
    )
    enriched = enrich_evidence_ref_location(ref, (parsed,))

    assert enriched.source_type == "txt"
    assert enriched.locator_confidence == "exact"
    assert enriched.locator["line_start"] == 2
    assert enriched.locator["line_end"] == 2
    assert enriched.anchor_id.endswith("#line-2")
    assert enriched.context == (
        "开场白\n客户说每年不能超过80万\n销售回应可以做预算释放"
    )


def test_enrich_ref_uses_model_provided_line_span_when_quote_is_paraphrased(
    tmp_path,
) -> None:
    source = tmp_path / "第1节.track-0.txt"
    source.write_text(
        "开场白\n"
        "他这90万的赔款\n"
        "在一周以来\n"
        "全部到账\n"
        "客户很认可\n",
        encoding="utf-8",
    )
    parsed = parse_source_file(source)
    ref = EvidenceRef(
        filename="第1节.track-0.txt",
        quote="他这90万的赔款在一周以内全部到账",
        locator={"line_start": 2, "line_end": 4},
    )

    enriched = enrich_evidence_ref_location(ref, (parsed,))

    assert enriched.source_path == str(source)
    assert enriched.locator["line_start"] == 2
    assert enriched.locator["line_end"] == 4
    assert enriched.locator["char_start"] == len("开场白\n")
    assert enriched.locator_confidence == "validated_span"
    assert enriched.source_excerpt == "他这90万的赔款\n在一周以来\n全部到账"
    assert enriched.context == (
        "开场白\n他这90万的赔款\n在一周以来\n全部到账\n客户很认可"
    )


def test_enrich_ref_marks_unmatched_when_line_span_and_quote_do_not_resolve(
    tmp_path,
) -> None:
    source = tmp_path / "第1节.track-0.txt"
    source.write_text("开场白\n客户关注预算\n", encoding="utf-8")
    parsed = parse_source_file(source)
    ref = EvidenceRef(
        filename="第1节.track-0.txt",
        quote="完全不存在的归纳句",
        locator={"line_start": 99, "line_end": 100},
    )

    enriched = enrich_evidence_ref_location(ref, (parsed,))

    assert enriched.source_type == "txt"
    assert enriched.source_path == str(source)
    assert enriched.locator == {}
    assert enriched.locator_confidence == "unmatched"
    assert enriched.locator_error == "line_span_invalid_and_quote_not_found"


def test_load_case_sections_groups_supported_files_by_section_dir(tmp_path) -> None:
    case_dir = tmp_path / "案例A"
    section_dir = case_dir / "第1节"
    section_dir.mkdir(parents=True)
    (section_dir / "讲义.txt").write_text("保单整理能发现保障缺口", encoding="utf-8")
    (section_dir / "视频.mp4").write_text("ignored", encoding="utf-8")

    sections = load_case_sections(case_dir)

    assert len(sections) == 1
    assert sections[0].case_name == "案例A"
    assert sections[0].section_name == "第1节"
    assert [source.filename for source in sections[0].sources] == ["讲义.txt"]
    assert sections[0].skipped_files == ("视频.mp4",)


def test_load_case_sections_deduplicates_identical_source_text(tmp_path) -> None:
    case_dir = tmp_path / "案例A"
    section_dir = case_dir / "第1节"
    section_dir.mkdir(parents=True)
    (section_dir / "讲义.txt").write_text("同一份素材", encoding="utf-8")
    (section_dir / "讲义(1).txt").write_text("同一份素材", encoding="utf-8")

    sections = load_case_sections(case_dir)

    assert [source.filename for source in sections[0].sources] == ["讲义.txt"]


def test_parse_docx_source_extracts_embedded_images(tmp_path) -> None:
    from docx import Document

    image_path = tmp_path / "chart.png"
    image_path.write_bytes(_png_bytes())
    docx_path = tmp_path / "讲义.docx"
    document = Document()
    document.add_paragraph("请看这张保单整理表")
    document.add_picture(str(image_path))
    document.save(str(docx_path))

    parsed = parse_source_file(docx_path)

    assert parsed.source_type == "docx"
    assert parsed.images
    assert parsed.images[0].media_type == "image/png"
    assert parsed.images[0].data == image_path.read_bytes()


def test_parse_pptx_source_extracts_embedded_images(tmp_path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    image_path = tmp_path / "chart.png"
    image_path.write_bytes(_png_bytes())
    pptx_path = tmp_path / "讲义.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    presentation.save(str(pptx_path))

    parsed = parse_source_file(pptx_path)

    assert parsed.source_type == "pptx"
    assert parsed.images
    assert parsed.images[0].media_type == "image/png"
    assert parsed.images[0].data == image_path.read_bytes()


def test_parse_pdf_source_extracts_embedded_images(tmp_path) -> None:
    pdf_path = tmp_path / "讲义.pdf"
    Image.open(BytesIO(_png_bytes())).save(pdf_path, "PDF")

    parsed = parse_source_file(pdf_path)

    assert parsed.source_type == "pdf"
    assert parsed.images
    assert parsed.images[0].media_type.startswith("image/")
    assert parsed.images[0].data
