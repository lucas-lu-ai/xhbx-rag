from __future__ import annotations

import re
import hashlib
import mimetypes
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

from .models import EvidenceRef


SUPPORTED_EXTENSIONS = {
    ".docx": "docx",
    ".pptx": "pptx",
    ".pdf": "pdf",
    ".txt": "txt",
}
MEDIA_EXTENSIONS = {".mp3", ".mp4", ".wav", ".m4a", ".mov", ".ts"}
_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
_PAGE_RE = re.compile(r"^##\s*第\s*(\d+)\s*页\s*$")
_ELLIPSIS_RE = re.compile(r"\.{3,}|…+")
_DOCX_HEADING_STYLE_RE = re.compile(r"(?:heading|标题)\s*(\d+)", re.IGNORECASE)


@dataclass(frozen=True)
class ParsedEmbeddedImage:
    image_id: str
    filename: str
    source_path: str
    media_type: str
    data: bytes
    locator: dict[str, object] = field(default_factory=dict)


@dataclass(frozen=True)
class ParsedSourceFile:
    source_id: str
    source_type: str
    filename: str
    source_path: str
    text: str
    images: tuple[ParsedEmbeddedImage, ...] = ()
    warnings: tuple[str, ...] = ()

    @property
    def is_empty(self) -> bool:
        return not self.text.strip()


@dataclass(frozen=True)
class SourceSection:
    case_name: str
    section_name: str
    section_dir: str
    sources: tuple[ParsedSourceFile, ...]
    skipped_files: tuple[str, ...] = ()

    @property
    def primary_text(self) -> str:
        return "\n\n".join(source.text for source in self.sources if source.text.strip())


def load_case_sections(case_dir: Path, case_name: str | None = None) -> list[SourceSection]:
    if not case_dir.is_dir():
        raise NotADirectoryError(f"案例目录不存在或不是目录: {case_dir}")

    resolved_case_name = case_name or case_dir.name
    sections: list[SourceSection] = []
    for directory in _candidate_section_dirs(case_dir):
        sources: list[ParsedSourceFile] = []
        skipped: list[str] = []
        for child in sorted(directory.iterdir()):
            if not child.is_file():
                continue
            suffix = child.suffix.lower()
            if suffix in SUPPORTED_EXTENSIONS:
                sources.append(parse_source_file(child, base_dir=case_dir.parent))
            elif suffix in MEDIA_EXTENSIONS:
                skipped.append(child.name)
        if not sources:
            continue
        sections.append(
            SourceSection(
                case_name=resolved_case_name,
                section_name=directory.name,
                section_dir=_relative_path(directory, case_dir.parent),
                sources=tuple(_deduplicate_sources(sources)),
                skipped_files=tuple(skipped),
            )
        )
    return sections


def parse_source_file(path: Path, *, base_dir: Path | None = None) -> ParsedSourceFile:
    source_type = SUPPORTED_EXTENSIONS.get(path.suffix.lower())
    if source_type is None:
        raise ValueError(f"不支持的素材文件类型: {path}")
    parser = _PARSERS[source_type]
    text, warnings, images = parser(path)
    text = _sanitize_surrogates(text)
    source_path = _relative_path(path, base_dir) if base_dir is not None else str(path)
    return ParsedSourceFile(
        source_id=f"{source_type}:{path.name}",
        source_type=source_type,
        filename=path.name,
        source_path=source_path,
        text=text,
        images=tuple(_with_relative_image_paths(images, source_path)),
        warnings=warnings,
    )


def enrich_evidence_ref_location(
    ref: EvidenceRef,
    sources: tuple[ParsedSourceFile, ...] | list[ParsedSourceFile],
) -> EvidenceRef:
    candidates = _candidate_sources(ref, sources)
    for source in candidates:
        span = _validated_line_span(ref, source)
        if span is not None:
            return _ref_with_resolved_span(
                ref,
                source,
                span[0],
                span[1],
                confidence="validated_span",
            )
        match = _find_quote_match(ref.quote, source.text)
        if match is None:
            continue
        start, end, confidence = match
        line_start = _line_number(source.text, start)
        line_end = _line_number(source.text, max(start, end - 1))
        return _ref_with_resolved_span(
            ref,
            source,
            line_start,
            line_end,
            confidence=confidence,
            char_start=start,
            char_end=end,
        )

    if candidates:
        error = (
            "line_span_invalid_and_quote_not_found"
            if _has_line_span(ref)
            else "quote_not_found"
        )
        return _fill_source_identity(ref, candidates[0]).model_copy(
            update={
                "locator": {},
                "locator_confidence": ref.locator_confidence or "unmatched",
                "locator_error": ref.locator_error or error,
            }
        )
    return ref


def enrich_evidence_refs(
    refs: list[EvidenceRef],
    sources: tuple[ParsedSourceFile, ...] | list[ParsedSourceFile],
) -> list[EvidenceRef]:
    return [enrich_evidence_ref_location(ref, sources) for ref in refs]


def _candidate_section_dirs(case_dir: Path) -> list[Path]:
    directories = [case_dir]
    directories.extend(path for path in sorted(case_dir.rglob("*")) if path.is_dir())
    return directories


def _relative_path(path: Path, base_dir: Path | None) -> str:
    if base_dir is None:
        return str(path)
    try:
        return str(path.resolve().relative_to(base_dir.resolve()))
    except ValueError:
        return str(path)


def _deduplicate_sources(sources: list[ParsedSourceFile]) -> list[ParsedSourceFile]:
    by_content: dict[tuple[str, str, tuple[str, ...]], ParsedSourceFile] = {}
    for source in sources:
        key = (
            source.source_type,
            hashlib.sha256(source.text.encode("utf-8")).hexdigest(),
            tuple(hashlib.sha256(image.data).hexdigest() for image in source.images),
        )
        existing = by_content.get(key)
        if existing is None or _source_preference_key(source) < _source_preference_key(
            existing
        ):
            by_content[key] = source
    return sorted(by_content.values(), key=lambda source: source.filename)


def _source_preference_key(source: ParsedSourceFile) -> tuple[int, int, str]:
    has_copy_suffix = 1 if re.search(r"\(\d+\)(?=\.[^.]+$)", source.filename) else 0
    is_empty = 1 if source.is_empty else 0
    return has_copy_suffix, is_empty, source.filename


def _sanitize_surrogates(text: str) -> str:
    """剔除孤立 UTF-16 代理字符（Office 文件中的数学符号等偶发产生），
    否则后续 utf-8 编码（hash/JSON/入库）会抛 UnicodeEncodeError。"""
    try:
        text.encode("utf-8")
        return text
    except UnicodeEncodeError:
        return text.encode("utf-8", errors="ignore").decode("utf-8")


def _parse_txt(path: Path) -> tuple[str, tuple[str, ...], tuple[ParsedEmbeddedImage, ...]]:
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            text = path.read_text(encoding="gbk")
        except Exception as exc:  # noqa: BLE001
            return "", (f"txt 解码失败: {exc}",), ()
    except Exception as exc:  # noqa: BLE001
        return "", (f"txt 读取失败: {exc}",), ()
    return text, () if text.strip() else ("txt 内容为空",), ()


def _parse_docx(path: Path) -> tuple[str, tuple[str, ...], tuple[ParsedEmbeddedImage, ...]]:
    try:
        from docx import Document
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError:
        return "", ("缺少 python-docx，无法解析 docx",), ()

    try:
        document = Document(str(path))
    except Exception as exc:  # noqa: BLE001
        return "", (f"docx 解析失败: {exc}",), ()

    blocks: list[str] = []
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            paragraph = Paragraph(child, document)
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = paragraph.style.name if paragraph.style else ""
            level = _heading_level(style_name)
            blocks.append(f"{'#' * level} {text}" if level else text)
        elif isinstance(child, CT_Tbl):
            table_text = _table_to_markdown(Table(child, document))
            if table_text:
                blocks.append(table_text)
    text = "\n\n".join(blocks)
    images = _extract_zip_images(path, prefix="word/media/", source_type="docx")
    return text, () if text.strip() else ("docx 内容为空",), images


def _parse_pptx(path: Path) -> tuple[str, tuple[str, ...], tuple[ParsedEmbeddedImage, ...]]:
    try:
        from pptx import Presentation
    except ImportError:
        return "", ("缺少 python-pptx，无法解析 pptx",), ()

    try:
        presentation = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        return "", (f"pptx 解析失败: {exc}",), ()

    blocks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    line = " | ".join(cell for cell in cells if cell)
                    if line:
                        parts.append(line)
        notes = _slide_notes_text(slide)
        if notes:
            parts.append(f"### 讲师备注\n{notes}")
        if parts:
            blocks.append(f"## 第 {index} 页\n" + "\n".join(parts))
    text = "\n\n".join(blocks)
    images = _extract_zip_images(path, prefix="ppt/media/", source_type="pptx")
    return text, () if text.strip() else ("pptx 无文本内容",), images


def _slide_notes_text(slide: Any) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    notes_frame = slide.notes_slide.notes_text_frame
    if notes_frame is None:
        return ""
    return notes_frame.text.strip()


def _parse_pdf(path: Path) -> tuple[str, tuple[str, ...], tuple[ParsedEmbeddedImage, ...]]:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "", ("缺少 pypdf，无法解析 pdf",), ()

    try:
        reader = PdfReader(str(path))
    except Exception as exc:  # noqa: BLE001
        return "", (f"pdf 解析失败: {exc}",), ()

    blocks: list[str] = []
    images: list[ParsedEmbeddedImage] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:  # noqa: BLE001
            text = ""
        if text:
            blocks.append(f"## 第 {index} 页\n{text}")
        images.extend(_extract_pdf_page_images(page, path=path, page_number=index))
    text = "\n\n".join(blocks)
    return text, () if text.strip() else ("pdf 无可提取文本（可能为扫描件，需 OCR）",), tuple(images)


_PARSERS: dict[
    str,
    Callable[[Path], tuple[str, tuple[str, ...], tuple[ParsedEmbeddedImage, ...]]],
] = {
    "txt": _parse_txt,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "pdf": _parse_pdf,
}


def _extract_zip_images(
    path: Path,
    *,
    prefix: str,
    source_type: str,
) -> tuple[ParsedEmbeddedImage, ...]:
    images: list[ParsedEmbeddedImage] = []
    try:
        with zipfile.ZipFile(path) as archive:
            media_names = [
                name
                for name in sorted(archive.namelist())
                if name.startswith(prefix) and not name.endswith("/")
            ]
            for index, name in enumerate(media_names, start=1):
                media_type = _image_media_type(name)
                if media_type is None:
                    continue
                images.append(
                    ParsedEmbeddedImage(
                        image_id=f"{source_type}:{path.name}:image-{index}",
                        filename=Path(name).name,
                        source_path=f"{path}::{name}",
                        media_type=media_type,
                        data=archive.read(name),
                        locator={"container": name},
                    )
                )
    except Exception:  # noqa: BLE001 - image extraction is best-effort
        return ()
    return tuple(images)


def _extract_pdf_page_images(
    page,
    *,
    path: Path,
    page_number: int,
) -> list[ParsedEmbeddedImage]:
    images: list[ParsedEmbeddedImage] = []
    try:
        page_images = list(getattr(page, "images", []))
    except Exception:  # noqa: BLE001
        return images
    for index, image in enumerate(page_images, start=1):
        data = getattr(image, "data", b"")
        if not data:
            continue
        filename = getattr(image, "name", "") or f"page-{page_number}-image-{index}"
        media_type = _image_media_type(filename) or "image/jpeg"
        images.append(
            ParsedEmbeddedImage(
                image_id=f"pdf:{path.name}:page-{page_number}-image-{index}",
                filename=filename,
                source_path=f"{path}::page-{page_number}::{filename}",
                media_type=media_type,
                data=data,
                locator={"page": page_number},
            )
        )
    return images


def _image_media_type(filename: str) -> str | None:
    media_type, _ = mimetypes.guess_type(filename)
    if media_type and media_type.startswith("image/"):
        return media_type
    suffix = Path(filename).suffix.lower()
    if suffix in {".jpg", ".jpeg"}:
        return "image/jpeg"
    if suffix == ".png":
        return "image/png"
    if suffix == ".gif":
        return "image/gif"
    if suffix == ".bmp":
        return "image/bmp"
    if suffix == ".webp":
        return "image/webp"
    if suffix in {".tif", ".tiff"}:
        return "image/tiff"
    return None


def _with_relative_image_paths(
    images: tuple[ParsedEmbeddedImage, ...],
    source_path: str,
) -> list[ParsedEmbeddedImage]:
    relocated: list[ParsedEmbeddedImage] = []
    for image in images:
        _, _, inner_path = image.source_path.partition("::")
        relocated.append(
            ParsedEmbeddedImage(
                image_id=image.image_id,
                filename=image.filename,
                source_path=f"{source_path}::{inner_path}" if inner_path else source_path,
                media_type=image.media_type,
                data=image.data,
                locator=image.locator,
            )
        )
    return relocated


def _heading_level(style_name: str | None) -> int | None:
    if not style_name:
        return None
    match = _DOCX_HEADING_STYLE_RE.search(style_name)
    return int(match.group(1)) if match else None


def _table_to_markdown(table) -> str:
    rows = [
        [cell.text.strip().replace("\n", " ") for cell in row.cells]
        for row in table.rows
    ]
    if not rows:
        return ""
    header = rows[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _candidate_sources(
    ref: EvidenceRef,
    sources: tuple[ParsedSourceFile, ...] | list[ParsedSourceFile],
) -> list[ParsedSourceFile]:
    if ref.source_path:
        matched = [
            source for source in sources if source.source_path == ref.source_path
        ]
        if matched:
            return matched
    if ref.filename:
        matched = [source for source in sources if source.filename == ref.filename]
        if matched:
            return matched
    if ref.source_id:
        matched = [source for source in sources if source.source_id == ref.source_id]
        if matched:
            return matched
    return list(sources)


def _fill_source_identity(ref: EvidenceRef, source: ParsedSourceFile) -> EvidenceRef:
    return ref.model_copy(
        update={
            "source_id": ref.source_id or source.source_id,
            "filename": ref.filename or source.filename,
            "source_type": ref.source_type or source.source_type,
            "source_path": ref.source_path or source.source_path,
        }
    )


def _has_line_span(ref: EvidenceRef) -> bool:
    locator = ref.locator or {}
    return "line_start" in locator or "line_end" in locator


def _validated_line_span(
    ref: EvidenceRef,
    source: ParsedSourceFile,
) -> tuple[int, int] | None:
    locator = ref.locator or {}
    try:
        line_start = int(locator.get("line_start"))
        line_end = int(locator.get("line_end"))
    except (TypeError, ValueError):
        return None
    if line_start < 1 or line_end < line_start:
        return None
    line_count = len(source.text.splitlines())
    if line_count == 0 or line_end > line_count:
        return None
    return line_start, line_end


def _ref_with_resolved_span(
    ref: EvidenceRef,
    source: ParsedSourceFile,
    line_start: int,
    line_end: int,
    *,
    confidence: str,
    char_start: int | None = None,
    char_end: int | None = None,
) -> EvidenceRef:
    start, end = (
        (char_start, char_end)
        if char_start is not None and char_end is not None
        else _line_span_offsets(source.text, line_start, line_end)
    )
    locator = dict(ref.locator)
    locator.update(
        {
            "line_start": line_start,
            "line_end": line_end,
            "char_start": start,
            "char_end": end,
        }
    )
    locator.update(_structured_locator(source, source.text, line_start))
    return ref.model_copy(
        update={
            "source_id": ref.source_id or source.source_id,
            "filename": ref.filename or source.filename,
            "source_type": ref.source_type or source.source_type,
            "source_path": ref.source_path or source.source_path,
            "context": _context_window(source.text, line_start, line_end),
            "source_excerpt": _line_span_excerpt(source.text, line_start, line_end),
            "locator": locator,
            "locator_confidence": ref.locator_confidence or confidence,
            "locator_error": "",
            "anchor_id": ref.anchor_id or f"{source.source_id}#line-{line_start}",
        }
    )


def _line_span_offsets(text: str, line_start: int, line_end: int) -> tuple[int, int]:
    lines = text.splitlines(keepends=True)
    start = sum(len(line) for line in lines[: line_start - 1])
    end = start + sum(len(line) for line in lines[line_start - 1 : line_end])
    while end > start and text[end - 1] in "\r\n":
        end -= 1
    return start, end


def _line_span_excerpt(text: str, line_start: int, line_end: int) -> str:
    lines = text.splitlines()
    return "\n".join(lines[line_start - 1 : line_end])


def _find_quote_match(quote: str, text: str) -> tuple[int, int, str] | None:
    stripped = quote.strip()
    if not stripped or not text:
        return None
    start = text.find(stripped)
    if start >= 0:
        return start, start + len(stripped), "exact"

    parts = [part.strip() for part in _ELLIPSIS_RE.split(stripped) if part.strip()]
    for part in sorted(parts, key=len, reverse=True):
        if len(part) < 2:
            continue
        start = text.find(part)
        if start >= 0:
            return start, start + len(part), "approximate"
    return None


def _line_number(text: str, offset: int) -> int:
    return text.count("\n", 0, max(offset, 0)) + 1


def _context_window(
    text: str,
    line_start: int,
    line_end: int,
    *,
    before: int = 3,
    after: int = 3,
) -> str:
    lines = text.splitlines()
    if not lines:
        return ""
    start = max(1, line_start - before)
    end = min(len(lines), line_end + after)
    return "\n".join(line.strip() for line in lines[start - 1 : end] if line.strip())


def _structured_locator(
    source: ParsedSourceFile,
    text: str,
    line_number: int,
) -> dict[str, object]:
    lines = text.splitlines()
    locator: dict[str, object] = {}
    if source.source_type in {"pptx", "pdf"}:
        page = _page_at_line(lines, line_number)
        if page is not None:
            locator["slide" if source.source_type == "pptx" else "page"] = page
    if source.source_type == "docx":
        heading_path = _heading_path_at_line(lines, line_number)
        if heading_path:
            locator["heading_path"] = heading_path
    return locator


def _page_at_line(lines: list[str], line_number: int) -> int | None:
    current: int | None = None
    for index, line in enumerate(lines, start=1):
        if index > line_number:
            break
        match = _PAGE_RE.match(line.strip())
        if match:
            current = int(match.group(1))
    return current


def _heading_path_at_line(lines: list[str], line_number: int) -> list[str]:
    stack: list[str | None] = []
    for index, line in enumerate(lines, start=1):
        if index > line_number:
            break
        match = _HEADING_RE.match(line)
        if not match:
            continue
        level = len(match.group(1))
        title = match.group(2).strip()
        if len(stack) < level:
            stack.extend([None] * (level - len(stack)))
        stack = stack[:level]
        stack[level - 1] = title
    return [item for item in stack if item]
